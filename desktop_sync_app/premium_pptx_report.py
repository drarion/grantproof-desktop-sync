from __future__ import annotations

import tempfile
from pathlib import Path
from typing import Any

from PIL import Image, ImageOps, ImageEnhance
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from premium_report import PremiumActivityReportBuilder, _safe_text

BLUE = '064890'
BLUE_2 = '0A57B5'
BLUE_3 = '1665C1'
ORANGE = 'F56D28'
DARK = '23324A'
MUTED = '53647B'
CARD_BORDER = 'DCE6F3'
CARD_BG = 'FFFFFF'
PALE_BLUE = 'F5F9FF'
SOFT_BLUE = 'F2F6FB'
WHITE = 'FFFFFF'

HERO_W = 3.74
HERO_H = 2.50
ANNEX_W = 5.66
ANNEX_H = 2.04
MAP_W = 3.78
MAP_H = 1.24
# Measured for the PPTX overview body box: 3.38in x 2.22in at 8.7pt, line spacing 1.03.
# The business-writing layer must stay under this budget so the value paragraph is never hidden or cut.
ACTIVITY_OVERVIEW_MAX_CHARS = 600
STORY_OVERVIEW_MAX_CHARS = 560


def _rgb(hex_color: str) -> RGBColor:
    clean = hex_color.replace('#', '')
    return RGBColor(int(clean[0:2], 16), int(clean[2:4], 16), int(clean[4:6], 16))


def _crop_cover(source: Path, output: Path, size: tuple[int, int], radius: int = 22) -> Path:
    im = Image.open(source).convert('RGB')
    tw, th = size
    ratio = tw / th
    if im.width / max(im.height, 1) > ratio:
        new_w = int(im.height * ratio)
        x = max(0, (im.width - new_w) // 2)
        im = im.crop((x, 0, x + new_w, im.height))
    else:
        new_h = int(im.width / ratio)
        y = max(0, int((im.height - new_h) * 0.35))
        im = im.crop((0, y, im.width, y + new_h))
    im = im.resize((tw, th), Image.Resampling.LANCZOS).convert('RGBA')
    mask = Image.new('L', im.size, 0)
    from PIL import ImageDraw
    d = ImageDraw.Draw(mask)
    d.rounded_rectangle((0, 0, tw - 1, th - 1), radius=radius, fill=255)
    im.putalpha(mask)
    im.save(output)
    return output


def _fit_contain(source: Path, output: Path, size: tuple[int, int], radius: int = 18, bg: str = '#FFFFFF') -> Path:
    im = Image.open(source).convert('RGBA')
    tw, th = size
    contain = ImageOps.contain(im, (tw, th), Image.Resampling.LANCZOS)
    canvas = Image.new('RGBA', (tw, th), bg)
    ox = (tw - contain.width) // 2
    oy = (th - contain.height) // 2
    canvas.alpha_composite(contain, (ox, oy))
    mask = Image.new('L', (tw, th), 0)
    from PIL import ImageDraw
    d = ImageDraw.Draw(mask)
    d.rounded_rectangle((0, 0, tw - 1, th - 1), radius=radius, fill=255)
    canvas.putalpha(mask)
    canvas.save(output)
    return output


class PremiumPptxActivityReportBuilder:
    """Editable PowerPoint premium activity report.

    The visual report is generated as native PowerPoint objects: editable text boxes,
    native shapes, independent photo/map/icon images, and no full-slide flattened image.
    """

    def __init__(self, base_folder: Path, default_org_name: str = '') -> None:
        self.base_folder = Path(base_folder)
        self.default_org_name = _safe_text(default_org_name)
        self.helper = PremiumActivityReportBuilder(base_folder, default_org_name=default_org_name)
        self._tmp_dir: Path | None = None

    def build(self, output: Path, project_code: str, project_name: str, records: list[Any], lang: str) -> None:
        lang = 'en' if str(lang).lower().startswith('en') else 'fr'
        if not records:
            return
        output.parent.mkdir(parents=True, exist_ok=True)
        data = self.helper._prepare_data(project_code, project_name, records, lang)
        self._apply_overview_text_budget(data, lang)
        with tempfile.TemporaryDirectory(prefix='grantproof_pptx_report_assets_') as tmp:
            self._tmp_dir = Path(tmp)
            self.helper._temp_dir = self._tmp_dir
            self.helper._icon_cache = {}
            data['hero_picture'] = self._prepare_hero(data.get('hero_image'), self._tmp_dir / 'hero_round.png')
            raw_map = self.helper._render_map_image(data, self._tmp_dir / 'map_raw.png', (1520, 500), lang)
            data['map_picture'] = self._prepare_map(raw_map, self._tmp_dir / 'map.png')
            data['header_country_picture'] = self.helper._render_country_silhouette(data, self._tmp_dir / 'country.png', (360, 230), for_header=True)

            prs = Presentation()
            prs.slide_width = Inches(13.333333)
            prs.slide_height = Inches(7.5)
            self._build_dashboard_slide(prs, data, lang)
            self._build_annex_slides(prs, data, lang)
            prs.save(output)
            self.helper._temp_dir = None
            self.helper._icon_cache = {}
            self._tmp_dir = None

    # ------------------------------------------------------------------ slide builders

    def _build_dashboard_slide(self, prs: Presentation, data: dict[str, Any], lang: str) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_background(slide, WHITE)
        W, H = 13.333333, 7.5

        # Header — mimics the provided institutional model.
        self._rect(slide, 0, 0, W, 1.18, BLUE, line=BLUE)
        title = data['title'].upper()
        long_title = len(title) > 60
        title_size = 18.8 if long_title else 22
        title_y = 0.18 if long_title else 0.28
        title_h = 0.50 if long_title else 0.33
        subtitle_y = 0.78 if long_title else 0.76
        self._text(slide, 0.30, title_y, 8.55, title_h, title, size=title_size, color=WHITE, bold=True, font='Arial Narrow')
        subtitle_prefix = 'Projet' if lang == 'fr' else 'Project'
        self._text(slide, 0.31, subtitle_y, 8.70, 0.20, f"{subtitle_prefix} : {data['project_name']} ({data['project_code']})", size=11.6 if long_title else 12.6, color=WHITE, bold=True)
        # Subtle radial-like darker edge using transparent-ish same-color shapes is not reliable in PowerPoint, keep flat UN/OCHA blue.
        self._image(slide, data['header_country_picture'], 9.62, 0.12, 1.22, 0.90)
        self._line(slide, 11.10, 0.16, 11.10, 1.02, WHITE, width=1.1, transparency=25)
        meta_x, meta_y = 11.42, 0.23
        meta = [('meta_org', data['org_name']), ('meta_pin', data['location']), ('meta_calendar', data['report_date'])]
        for i, (icon_key, value) in enumerate(meta):
            y = meta_y + i * 0.31
            self._image(slide, self.helper._white_icon(icon_key), meta_x, y + 0.015, 0.15, 0.15)
            self._text(slide, meta_x + 0.34, y - 0.01, 1.58, 0.20, self._truncate(value, 32), size=11, color=WHITE, bold=False)

        # Top row cards.
        y_top, top_h = 1.34, 4.05
        overview_x, overview_w = 0.23, 3.80
        kpi_x, kpi_w = 4.33, 4.42
        media_x, media_w = 9.08, 4.02
        self._card(slide, overview_x, y_top, overview_w, top_h, fill=WHITE, shadow=True)
        self._card(slide, kpi_x, y_top, kpi_w, top_h, fill=WHITE, line=WHITE, shadow=True)
        self._card(slide, media_x, y_top, media_w, top_h, fill=WHITE, line=WHITE, shadow=True)
        self._line(slide, 8.88, y_top + 0.05, 8.88, y_top + top_h - 0.05, CARD_BORDER, width=1.0)

        self._overview(slide, data, lang, overview_x, y_top, overview_w, top_h)
        self._kpis(slide, data, lang, kpi_x, y_top, kpi_w, top_h)
        self._hero_and_sectors(slide, data, lang, media_x, y_top, media_w, top_h)

        self._line(slide, 0.25, 5.58, 13.08, 5.58, CARD_BORDER, width=1.0)

        # Bottom cards.
        self._card(slide, 0.23, 5.70, 6.13, 1.58, fill=SOFT_BLUE, line='E4EBF5', shadow=True)
        self._card(slide, 6.55, 5.70, 6.55, 1.58, fill=SOFT_BLUE, line='E4EBF5', shadow=True)
        self._highlights(slide, data, lang, 0.23, 5.70, 6.13, 1.58)
        self._location(slide, data, lang, 6.55, 5.70, 6.55, 1.58)

    def _build_annex_slides(self, prs: Presentation, data: dict[str, Any], lang: str) -> None:
        images = [Path(p) for p in data.get('annex_images') or [] if Path(p).exists()]
        if not images:
            return
        for start in range(0, len(images), 4):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._set_background(slide, WHITE)
            self._rect(slide, 0, 0, 13.333333, 0.90, BLUE, line=BLUE)
            title = 'ANNEXE DOCUMENTAIRE' if lang == 'fr' else 'DOCUMENTARY ANNEX'
            self._text(slide, 0.35, 0.25, 5.0, 0.35, title, size=20, color=WHITE, bold=True, font='Arial Narrow')
            self._text(slide, 0.35, 0.59, 8.0, 0.22, data['title'].upper(), size=10, color=WHITE)
            self._line(slide, 0.0, 0.90, 13.333333, 0.90, 'D7E2F2', width=1.0)
            positions = [(0.50, 1.25), (6.95, 1.25), (0.50, 4.35), (6.95, 4.35)]
            for idx, img_path in enumerate(images[start:start+4]):
                x, y = positions[idx]
                self._card(slide, x, y, 5.90, 2.55, fill=WHITE)
                prepared = _crop_cover(img_path, self._tmp_dir / f'annex_{start+idx}.png', (1698, 612), radius=28)  # type: ignore[operator]
                self._image(slide, prepared, x + 0.12, y + 0.12, ANNEX_W, ANNEX_H)
                self._text(slide, x + 0.12, y + 2.24, 5.66, 0.20, self._annex_caption(data, start + idx, lang), size=9.4, color=DARK)

    # ------------------------------------------------------------------ sections

    def _overview(self, slide, data: dict[str, Any], lang: str, x: float, y: float, w: float, h: float) -> None:
        title = 'APERÇU DE L’ACTIVITÉ' if lang == 'fr' else 'ACTIVITY OVERVIEW'
        if data.get('report_mode') == 'story':
            title = 'APERÇU DE LA SUCCESS STORY' if lang == 'fr' else 'SUCCESS STORY OVERVIEW'
        self._section_title(slide, x + 0.22, y + 0.22, w - 0.44, title, 'model_section_overview')
        # Never truncate the value text in display. It has already been generated/re-written to the measured budget.
        summary = _safe_text(data.get('summary'))
        self._text(slide, x + 0.20, y + 0.73, w - 0.42, 2.22, summary, size=8.7, color=DARK, fit=False, line_spacing=1.03)
        self._line(slide, x + 0.20, y + 3.03, x + w - 0.20, y + 3.03, CARD_BORDER, width=0.7)
        rows = [
            ('model_info_location', 'Lieu' if lang == 'fr' else 'Location', data['location']),
            ('model_info_people', 'Public cible' if lang == 'fr' else 'Target group', data['target_group']),
            ('model_info_photo', ('Type de story' if lang == 'fr' else 'Story type') if data.get('report_mode') == 'story' else ('Type de preuve' if lang == 'fr' else 'Evidence type'), data['evidence_type']),
        ]
        for i, (icon, label, value) in enumerate(rows):
            yy = y + 3.13 + i * 0.30
            self._image(slide, self.helper._plain_icon(icon), x + 0.25, yy + 0.02, 0.22, 0.22)
            self._text(slide, x + 0.58, yy - 0.005, w - 0.80, 0.15, label, size=7.8, color=BLUE_2, bold=True)
            self._text(slide, x + 0.58, yy + 0.14, w - 0.80, 0.14, self._truncate(value, 62), size=7.4, color=DARK)
            if i < 2:
                self._line(slide, x + 0.20, yy + 0.28, x + w - 0.20, yy + 0.28, 'E7EEF8', width=0.5)

    def _kpis(self, slide, data: dict[str, Any], lang: str, x: float, y: float, w: float, h: float) -> None:
        self._section_title(slide, x + 0.06, y + 0.22, w - 0.12, 'CHIFFRES CLÉS' if lang == 'fr' else 'KEY FIGURES', 'model_section_kpis', no_circle=True)
        kpis = self._normalized_kpis(data, lang)
        card_w, card_h = 1.34, 2.95
        gap = 0.16
        start_x = x + (w - (3*card_w + 2*gap)) / 2
        for i, kpi in enumerate(kpis[:3]):
            cx = start_x + i * (card_w + gap)
            self._card(slide, cx, y + 1.02, card_w, card_h, fill=WHITE, line='E8EDF5', shadow=True)
            self._oval(slide, cx + 0.30, y + 1.47, 0.76, 0.76, fill='F1F6FE', line='F1F6FE')
            self._image(slide, self.helper._plain_icon(self.helper._kpi_icon_key(kpi)), cx + 0.47, y + 1.64, 0.42, 0.42)
            self._text(slide, cx + 0.12, y + 2.43, card_w - 0.24, 0.44, self._truncate(str(kpi.get('value', '')), 8), size=25, color=BLUE_2, bold=True, align='center')
            self._line(slide, cx + 0.57, y + 3.04, cx + 0.77, y + 3.04, ORANGE, width=1.2)
            label = self._truncate(str(kpi.get('label', '')), 24)
            self._text(slide, cx + 0.11, y + 3.20, card_w - 0.22, 0.52, label, size=10.0, color=DARK, align='center')

    def _hero_and_sectors(self, slide, data: dict[str, Any], lang: str, x: float, y: float, w: float, h: float) -> None:
        self._image(slide, data['hero_picture'], x + 0.14, y + 0.12, HERO_W, HERO_H)
        self._section_title(slide, x + 0.15, y + 2.92, w - 0.30, 'ALIGNEMENT SECTORIEL' if lang == 'fr' else 'SECTOR ALIGNMENT', 'model_section_sector')
        sectors = self._normalized_sectors(data)
        sx1, sx2 = x + 0.52, x + 2.32
        sy = y + 3.47
        for idx, (sx, sector) in enumerate([(sx1, sectors[0]), (sx2, sectors[1])]):
            icon_key = 'food_security' if sector.get('key') in {'food_security', 'food'} else ('agriculture' if sector.get('key') == 'agriculture' else sector.get('icon', 'coordination'))
            self._image(slide, self.helper._sector_icon(icon_key), sx, sy, 0.56, 0.56)
            self._text(slide, sx + 0.66, sy + 0.10, 1.05, 0.36, sector.get('label', ''), size=9.7, color=MUTED, bold=True if idx == 0 else False)
        self._line(slide, x + 2.03, sy - 0.02, x + 2.03, sy + 0.74, CARD_BORDER, width=0.9)

    def _highlights(self, slide, data: dict[str, Any], lang: str, x: float, y: float, w: float, h: float) -> None:
        self._section_title(slide, x + 0.22, y + 0.14, 4.3, 'FAITS SAILLANTS' if lang == 'fr' else 'HIGHLIGHTS', 'model_section_highlights')
        for i, item in enumerate(data['highlights'][:3]):
            yy = y + 0.67 + i * 0.30
            self._text(slide, x + 0.28, yy + 0.01, 0.15, 0.15, '•', size=14, color=BLUE_2, bold=True)
            self._text(slide, x + 0.56, yy, w - 0.80, 0.27, self._truncate(item, 118), size=8.9, color=DARK, fit=False)

    def _location(self, slide, data: dict[str, Any], lang: str, x: float, y: float, w: float, h: float) -> None:
        self._section_title(slide, x + 0.24, y + 0.14, 2.05, 'LOCALISATION' if lang == 'fr' else 'LOCATION', 'model_section_location')
        self._text(slide, x + 0.30, y + 0.98, 1.85, 0.20, 'Localisation de l’activité' if lang == 'fr' else 'Activity location', size=8.4, color=BLUE_2, bold=True)
        self._line(slide, x + 0.30, y + 1.20, x + 0.78, y + 1.20, ORANGE, width=1.2)
        if data.get('gps_point'):
            lat, lon, _ = data['gps_point']
            gps_text = f'{lat:.4f}, {lon:.4f}'
        else:
            gps_text = '—'
        self._text(slide, x + 0.30, y + 1.30, 1.92, 0.18, self._truncate(data['location'], 42), size=8.3, color=DARK)
        self._text(slide, x + 0.30, y + 1.43, 1.92, 0.18, f'GPS : {gps_text}' if gps_text != '—' else '', size=7.4, color=MUTED)
        self._image(slide, data['map_picture'], x + 2.55, y + 0.10, MAP_W, MAP_H)
        self._text(slide, x + 2.55, y + 1.36, 3.78, 0.18, data.get('map_location_label') or data['location'], size=8.6, color=BLUE_2, bold=True, align='center')

    # ------------------------------------------------------------------ text budget helpers

    def _apply_overview_text_budget(self, data: dict[str, Any], lang: str) -> None:
        budget = STORY_OVERVIEW_MAX_CHARS if data.get('report_mode') == 'story' else ACTIVITY_OVERVIEW_MAX_CHARS
        data['activity_overview_max_chars'] = budget
        summary = _safe_text(data.get('summary'))
        if len(summary) <= budget:
            return
        data['summary'] = self._rewrite_overview_with_budget(data, lang, budget)

    def _rewrite_overview_with_budget(self, data: dict[str, Any], lang: str, budget: int) -> str:
        mode = data.get('report_mode')
        if mode == 'story':
            if lang == 'fr':
                candidates = [
                    f"Cette success story met en lumière {data.get('target_group')} à {data.get('location')} dans le cadre du projet {data.get('project_name')}.",
                    f"Le récit documente les changements observés et l’expérience vécue autour de {str(data.get('activity', '')).lower()}.",
                    "Les éléments collectés apportent une lecture qualitative utile pour le reporting bailleur, la redevabilité et la valorisation humaine de l’intervention.",
                ]
            else:
                candidates = [
                    f"This success story highlights {data.get('target_group')} in {data.get('location')} under the project {data.get('project_name')}.",
                    f"The story documents observed change and lived experience linked to {str(data.get('activity', '')).lower()}.",
                    "The collected material provides qualitative evidence for donor reporting, accountability, and human-centered communication.",
                ]
        else:
            indicator_phrase = self._kpi_phrase(data, lang)
            if lang == 'fr':
                candidates = [
                    f"Cette activité de {str(data.get('activity', '')).lower()} a été documentée à {data.get('location')} dans le cadre du projet {data.get('project_name')}.",
                    f"Les données disponibles mettent en avant {indicator_phrase} et donnent une lecture claire du volume d’appui réalisé." if indicator_phrase else "Les données disponibles donnent une lecture claire du volume d’appui réalisé.",
                    "Les éléments collectés permettent d’illustrer la mise en œuvre, de soutenir le reporting bailleur et d’éclairer les décisions de suivi sans aller au-delà des informations disponibles.",
                ]
            else:
                candidates = [
                    f"This {str(data.get('activity', '')).lower()} activity was documented in {data.get('location')} under the project {data.get('project_name')}.",
                    f"Available data highlights {indicator_phrase} and provides a clear reading of the operational support delivered." if indicator_phrase else "Available data provides a clear reading of the operational support delivered.",
                    "The collected elements support donor reporting, illustrate implementation progress, and inform follow-up decisions without going beyond the available information.",
                ]
        text = ''
        for sentence in candidates:
            candidate = (text + ' ' + sentence).strip() if text else sentence
            if len(candidate) <= budget:
                text = candidate
        if text:
            return text
        fallback = candidates[0]
        return fallback[:max(0, budget - 1)].rstrip() + '…'

    def _kpi_phrase(self, data: dict[str, Any], lang: str) -> str:
        kpis = [item for item in self._normalized_kpis(data, lang) if item.get('value') not in {'—', ''}]
        focus = [f"{item.get('value')} {str(item.get('label', '')).lower()}" for item in kpis[:3]]
        if not focus:
            return ''
        if len(focus) == 1:
            return focus[0]
        joiner = ' et ' if lang == 'fr' else ' and '
        return ', '.join(focus[:-1]) + joiner + focus[-1]

    # ------------------------------------------------------------------ helper methods

    def _prepare_hero(self, source: Path | None, output: Path) -> Path:
        size = (1496, 1000)
        if source and Path(source).exists():
            return _crop_cover(Path(source), output, size, radius=30)
        return self.helper._prepare_photo(source, output, size)  # type: ignore[arg-type]

    def _prepare_map(self, source: Path, output: Path) -> Path:
        im = Image.open(source).convert('RGBA')
        rgb = im.convert('RGB')
        rgb = ImageEnhance.Contrast(rgb).enhance(1.18)
        rgb = ImageEnhance.Sharpness(rgb).enhance(1.12)
        rgb = ImageEnhance.Color(rgb).enhance(0.92)
        tmp = output.with_name(output.stem + '_enhanced.png')
        rgb.save(tmp)
        return _fit_contain(tmp, output, (1512, 496), radius=18, bg='#FFFFFF')

    def _set_background(self, slide, color: str) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = _rgb(color)

    def _rect(self, slide, x, y, w, h, fill: str, line: str | None = None, radius: bool = False):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
        try:
            if radius and len(shp.adjustments):
                shp.adjustments[0] = 0.08
        except Exception:
            pass
        shp.fill.solid(); shp.fill.fore_color.rgb = _rgb(fill)
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = _rgb(line)
            shp.line.width = Pt(0.7)
        return shp

    def _card(self, slide, x, y, w, h, fill=WHITE, line=CARD_BORDER, shadow=False):
        shp = self._rect(slide, x, y, w, h, fill, line=line, radius=True)
        if shadow:
            try:
                shp.shadow.inherit = False
                shp.shadow.blur_radius = Pt(6)
                shp.shadow.distance = Pt(1.5)
                shp.shadow.transparency = 78
            except Exception:
                pass
        return shp

    def _oval(self, slide, x, y, w, h, fill: str, line: str | None = None):
        shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid(); shp.fill.fore_color.rgb = _rgb(fill)
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = _rgb(line)
            shp.line.width = Pt(0.5)
        return shp

    def _line(self, slide, x1, y1, x2, y2, color: str, width: float = 1.0, transparency: int | None = None):
        line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.line.color.rgb = _rgb(color)
        line.line.width = Pt(width)
        if transparency is not None:
            try:
                line.line.fill.transparency = transparency
            except Exception:
                pass
        return line

    def _text(self, slide, x, y, w, h, text: str, size=10, color=DARK, bold=False, font='Arial', align='left', fit=True, line_spacing=None):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        box.text_frame.clear()
        box.text_frame.margin_left = 0
        box.text_frame.margin_right = 0
        box.text_frame.margin_top = 0
        box.text_frame.margin_bottom = 0
        box.text_frame.word_wrap = True
        box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
        p = box.text_frame.paragraphs[0]
        if align == 'center':
            p.alignment = PP_ALIGN.CENTER
        elif align == 'right':
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = str(text)
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(color)
        return box

    def _image(self, slide, path: Path | str, x, y, w, h):
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))

    def _section_title(self, slide, x: float, y: float, w: float, title: str, icon_key: str, no_circle: bool = False) -> None:
        if no_circle:
            self._image(slide, self.helper._plain_icon(icon_key), x, y + 0.02, 0.36, 0.36)
        else:
            self._oval(slide, x, y + 0.00, 0.42, 0.42, fill=BLUE, line=BLUE)
            self._image(slide, self.helper._white_icon(icon_key), x + 0.105, y + 0.105, 0.21, 0.21)
        self._text(slide, x + 0.58, y + 0.07, w - 0.58, 0.27, title, size=12.5, color=BLUE_2, bold=True, font='Arial Narrow')
        self._line(slide, x + 0.58, y + 0.43, x + w, y + 0.43, BLUE_3, width=1.0, transparency=6)

    def _normalized_kpis(self, data: dict[str, Any], lang: str) -> list[dict[str, str]]:
        kpis = list(data.get('kpis') or [])
        if 'media' not in {item.get('kind') for item in kpis}:
            kpis.append({'kind': 'media', 'value': str(data.get('media_count', 0)), 'label': 'Médias liés' if lang == 'fr' else 'Linked media', 'icon': 'media'})
        if 'evidence' not in {item.get('kind') for item in kpis}:
            kpis.append({'kind': 'evidence', 'value': str(data.get('evidence_count', 1)), 'label': 'Preuve consolidée' if lang == 'fr' else 'Consolidated evidence', 'icon': 'check'})
        preferred = ['beneficiary', 'beneficiaries', 'people', 'farmer', 'kit', 'kits', 'group', 'groups', 'media', 'evidence']
        def rank(item):
            kind = _safe_text(item.get('kind')).lower()
            return (preferred.index(kind) if kind in preferred else 99, -len(_safe_text(item.get('value'))))
        result = []
        seen = set()
        for item in sorted(kpis, key=rank):
            sig = (item.get('kind'), item.get('value'))
            if sig in seen:
                continue
            result.append(item)
            seen.add(sig)
            if len(result) == 3:
                break
        while len(result) < 3:
            result.append({'kind': 'evidence', 'value': '—', 'label': 'Indicateur' if lang == 'fr' else 'Indicator', 'icon': 'check'})
        return result

    def _normalized_sectors(self, data: dict[str, Any]) -> list[dict[str, str]]:
        raw_sectors = list(data.get('sectors') or [])
        lang_is_fr = str(data.get('title', '')).lower().startswith(('rapport', 'success story :'))
        defaults = [
            {'key': 'food_security', 'label': 'Sécurité alimentaire' if lang_is_fr else 'Food security', 'icon': 'food_security'},
            {'key': 'agriculture', 'label': 'Agriculture', 'icon': 'agriculture'},
        ]
        # For agriculture / food-security activities, keep the institutional pair used in the reference model.
        keys = {s.get('key') for s in raw_sectors}
        if {'agriculture', 'food_security', 'food'} & keys:
            sectors: list[dict[str, str]] = []
            for default in defaults:
                matching = next((s for s in raw_sectors if s.get('key') in ({'food_security', 'food'} if default['key'] == 'food_security' else {default['key']})), None)
                sectors.append(matching or default)
            return sectors[:2]
        sectors = sorted(raw_sectors[:2], key=lambda s: 0 if s.get('key') in {'food_security','food'} else (1 if s.get('key') == 'agriculture' else 2))
        seen = {s.get('key') for s in sectors}
        for default in defaults:
            if len(sectors) >= 2:
                break
            if default['key'] not in seen:
                sectors.append(default)
                seen.add(default['key'])
        return sectors[:2]

    def _annex_caption(self, data: dict[str, Any], idx: int, lang: str) -> str:
        return f"{data['activity']} | {data['location']} | {data['report_date']}"

    def _truncate(self, text: Any, limit: int) -> str:
        clean = ' '.join(_safe_text(text).split())
        return clean if len(clean) <= limit else clean[:limit - 1].rstrip() + '…'
